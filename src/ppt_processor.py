# --- ppt_processor.py ---
from pptx import Presentation
import json
import zipfile
import os
import shutil
from pathlib import Path
import subprocess
import re
import time
from pdf2image import convert_from_path

def extract_text_from_shape(shape):
    """Extracts text from a shape, handling different shape types."""
    text = ""
    if shape.has_text_frame:
        text = shape.text_frame.text
    elif shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                text += cell.text_frame.text + "\t"
            text += "\n"
    return text.strip()

def pptx_to_json(filepath):
    """Converts a .pptx file to a JSON representation."""
    try:
        prs = Presentation(filepath)
        presentation_data = {
            "filename": os.path.basename(filepath),
            "slides": []
        }
        for i, slide in enumerate(prs.slides):
            slide_data = {
                "slide_number": i + 1,
                "shapes": [],
                "notes": ""
            }
            for shape in slide.shapes:
                shape_info = {
                    "name": shape.name,
                    "type": str(shape.shape_type),
                    "text": extract_text_from_shape(shape),
                    "left": shape.left.pt if hasattr(shape, 'left') and shape.left is not None else None,
                    "top": shape.top.pt if hasattr(shape, 'top') and shape.top is not None else None,
                    "width": shape.width.pt if hasattr(shape, 'width') and shape.width is not None else None,
                    "height": shape.height.pt if hasattr(shape, 'height') and shape.height is not None else None,
                }
                slide_data["shapes"].append(shape_info)
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                slide_data["notes"] = text_frame.text.strip()
            presentation_data["slides"].append(slide_data)
        return presentation_data
    except Exception as e:
        print(f"Error converting {filepath} to JSON: {e}")
        raise

def extract_xml_from_pptx(pptx_filepath, output_folder):
    """
    Extracts all constituent XML files from a .pptx file.
    Returns a list of full paths to the extracted XML files.
    """
    extracted_files_paths = []
    try:
        Path(output_folder).mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(pptx_filepath, 'r') as pptx_zip:
            for member_info in pptx_zip.infolist():
                member_name = member_info.filename
                if not member_info.is_dir() and member_name.endswith(('.xml', '.rels')):
                    target_path = os.path.join(output_folder, member_name)
                    os.makedirs(os.path.dirname(target_path), exist_ok=True)
                    with pptx_zip.open(member_name) as source, open(target_path, "wb") as target:
                        shutil.copyfileobj(source, target)
                    extracted_files_paths.append(target_path)
        return extracted_files_paths
    except Exception as e:
        print(f"Error extracting XML from {pptx_filepath}: {e}")
        raise

def create_modified_pptx(original_pptx_path, modified_xml_map, output_pptx_path):
    """
    Creates a new .pptx file by taking an original .pptx, and replacing
    specified internal XML files with new content.
    """
    temp_output_pptx_path = output_pptx_path + ".tmp"
    try:
        # Create the directory for the output file if it doesn't exist.
        os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
        with zipfile.ZipFile(original_pptx_path, 'r') as zin:
            with zipfile.ZipFile(temp_output_pptx_path, 'w', zipfile.ZIP_DEFLATED) as zout:
                for item in zin.infolist():
                    item_name_normalized = item.filename.replace("\\", "/")
                    if item_name_normalized in modified_xml_map:
                        new_content = modified_xml_map[item_name_normalized]
                        zout.writestr(item, new_content.encode('utf-8'))
                    else:
                        buffer = zin.read(item.filename)
                        zout.writestr(item, buffer)
        os.replace(temp_output_pptx_path, output_pptx_path)
        print(f"Modified PPTX successfully created at: {output_pptx_path}")
        return True
    except Exception as e:
        print(f"Error creating modified PPTX at {output_pptx_path}: {e}")
        if os.path.exists(temp_output_pptx_path):
            os.remove(temp_output_pptx_path)
        return False

def _find_soffice_command():
    """Finds a working LibreOffice/OpenOffice command."""
    soffice_commands = ['libreoffice', 'soffice']
    for cmd in soffice_commands:
        if shutil.which(cmd):
            try:
                subprocess.run([cmd, '--version'], capture_output=True, check=True, timeout=10)
                print(f"Found working soffice command: {cmd}")
                return cmd
            except Exception as e:
                print(f"Soffice command '{cmd}' not working or timed out: {e}")
    return None

def _convert_pptx_to_pdf(pptx_filepath, output_folder, soffice_cmd):
    """Converts a PPTX to a single PDF file using an isolated user profile for stability."""
    # Create a unique, temporary profile directory for this specific conversion process
    temp_profile_dir = Path(output_folder) / f"lo_profile_{os.getpid()}_{time.time_ns()}"
    temp_profile_dir.mkdir(parents=True, exist_ok=True)
    
    pdf_path = Path(output_folder) / (Path(pptx_filepath).stem + ".pdf")
    
    for attempt in range(2): # Retry mechanism
        try:
            command_args = [
                soffice_cmd,
                # ** FIX: Isolate LibreOffice instance to prevent parallel conflicts **
                f"-env:UserInstallation=file://{os.path.abspath(temp_profile_dir)}",
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', output_folder,
                pptx_filepath
            ]
            subprocess.run(command_args, capture_output=True, text=True, timeout=120, check=True)
            
            if pdf_path.exists():
                shutil.rmtree(temp_profile_dir, ignore_errors=True)
                return str(pdf_path)
            
            print(f"Attempt {attempt + 1}: PDF not found for {Path(pptx_filepath).name}. Retrying...")
            time.sleep(1)

        except subprocess.CalledProcessError as e:
            print(f"Attempt {attempt + 1}: Soffice error for {Path(pptx_filepath).name}. STDERR: {e.stderr.strip()}")
            time.sleep(1)
        except Exception as e:
            print(f"Attempt {attempt + 1}: Unexpected error during PDF conversion: {e}")
            break
    
    shutil.rmtree(temp_profile_dir, ignore_errors=True)
    print(f"Failed to convert {Path(pptx_filepath).name} to PDF after all attempts.")
    return None


def _convert_pdf_to_images(pdf_filepath, output_folder):
    """Converts a PDF file's pages to PNG images."""
    print(f"Converting PDF {pdf_filepath} to images...")
    try:
        images = convert_from_path(
            pdf_filepath,
            output_folder=output_folder,
            fmt='png',
            output_file='slide-',
            paths_only=True
        )
        print(f"Successfully converted PDF to {len(images)} images.")
        return sorted(images)
    except Exception as e:
        print(f"An error occurred converting PDF to images: {e}")
        print("Please ensure 'poppler' is installed on your system.")
        print("On macOS: 'brew install poppler'")
        print("On Debian/Ubuntu: 'sudo apt-get install poppler-utils'")
        return []



def export_slides_to_images(pptx_filepath, output_folder):
    """
    Robustly converts each slide of a .pptx file to a .png image by first
    converting to PDF, then splitting the PDF into images.
    """
    abs_pptx_filepath = os.path.abspath(pptx_filepath)
    abs_output_folder = os.path.abspath(output_folder)
    Path(abs_output_folder).mkdir(parents=True, exist_ok=True)

    soffice_cmd = _find_soffice_command()
    if not soffice_cmd:
        print("Error: LibreOffice command not found. Cannot proceed with image conversion.")
        return []

    pdf_path = _convert_pptx_to_pdf(abs_pptx_filepath, abs_output_folder, soffice_cmd)
    
    if not pdf_path:
        return []

    image_paths = _convert_pdf_to_images(pdf_path, abs_output_folder)

    try:
        os.remove(pdf_path)
        print(f"Cleaned up intermediate PDF: {pdf_path}")
    except OSError as e:
        print(f"Warning: Could not remove intermediate PDF {pdf_path}: {e}")

    return image_paths

def extract_specific_xml_from_pptx(pptx_filepath, xml_filename):
    """
    Extracts the content of a single specified XML file from a .pptx file.
    
    Args:
        pptx_filepath (str): Path to the .pptx file.
        xml_filename (str): The internal path to the XML file (e.g., 'ppt/slides/slide1.xml').
        
    Returns:
        str: The content of the XML file as a string, or None if not found.
    """
    try:
        with zipfile.ZipFile(pptx_filepath, 'r') as pptx_zip:
            # Normalize filename for matching
            xml_filename_normalized = xml_filename.replace("\\", "/")
            if xml_filename_normalized in pptx_zip.namelist():
                with pptx_zip.open(xml_filename_normalized) as xml_file:
                    return xml_file.read().decode('utf-8')
            return None
    except Exception as e:
        print(f"Error extracting specific XML '{xml_filename}' from {pptx_filepath}: {e}")
        return None